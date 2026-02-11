# -*- coding: utf-8 -*-
"""
PowerPointファイル用ツール

AIエージェントがPowerPointファイルを理解するための軽量ツール。
3つの機能を提供:
1. get_presentation_info: プレゼンテーション構造と基本情報を取得
2. get_slides_content: 指定スライドの内容を取得（範囲指定可能）
3. search_presentation: プレゼンテーション全体からキーワード検索
"""

import io
from typing import Any, TypedDict

import structlog

from app.services.workspace.file_tools.utils import (
    build_search_pattern,
    create_context_snippet,
    file_tool_handler,
    format_tool_error,
    format_tool_success,
    normalize_text,
)

logger = structlog.get_logger(__name__)


# =============================================================================
# Type Definitions
# =============================================================================

class SlideInfo(TypedDict):
    """スライド情報の型定義"""
    number: int
    title: str
    text_count: int
    image_count: int
    table_count: int
    chart_count: int
    char_count: int
    has_notes: bool
    notes_length: int


class PresentationInfo(TypedDict):
    """プレゼンテーション情報の型定義"""
    filename: str
    total_slides: int
    total_characters: int
    slides: list[SlideInfo]


class SlideContent(TypedDict):
    """スライドコンテンツの型定義"""
    number: int
    title: str
    text_content: list[str]
    table_content: list[str]
    notes: str | None


class SlidesResult(TypedDict):
    """スライド取得結果の型定義"""
    filename: str
    requested_slides: str
    total_slides: int
    returned_slides: int
    start_slide: int
    end_slide: int
    has_more: bool
    slides: list[SlideContent]


class SearchHit(TypedDict):
    """検索ヒットの型定義"""
    slide_number: int
    slide_title: str
    location_type: str  # "text" | "table" | "notes" | "title"
    text: str
    context: str


class SearchResult(TypedDict):
    """検索結果の型定義"""
    query: str
    total_hits: int
    hits: list[SearchHit]


# =============================================================================
# Constants
# =============================================================================

DEFAULT_MAX_SLIDES = 10  # デフォルトの最大取得スライド数


# =============================================================================
# Internal Utilities
# =============================================================================

def _load_presentation_from_bytes(content: bytes):
    """バイトデータからプレゼンテーションを読み込む"""
    from pptx import Presentation
    return Presentation(io.BytesIO(content))


def _get_slide_title(slide) -> str:
    """スライドのタイトルを取得"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            if hasattr(shape, "is_placeholder") and shape.placeholder_format:
                if shape.placeholder_format.type == 1:  # TITLE
                    return normalize_text(shape.text_frame.text).strip()[:100]
            elif shape.text_frame.text.strip():
                # 最初のテキストをタイトル候補として
                return normalize_text(shape.text_frame.text).strip()[:50]
    return "(タイトルなし)"


def _parse_slides(slides_spec: str, total_slides: int, max_slides: int) -> list[int]:
    """
    スライド指定を解析

    Args:
        slides_spec: "1-5" または "1,3,5" 形式
        total_slides: 総スライド数
        max_slides: 最大取得スライド数

    Returns:
        スライド番号のリスト（1始まり）
    """
    slide_numbers = []

    parts = slides_spec.replace(" ", "").split(",")

    for part in parts:
        if "-" in part:
            try:
                start, end = part.split("-")
                start_num = int(start)
                end_num = int(end)
                slide_numbers.extend(range(start_num, min(end_num + 1, total_slides + 1)))
            except ValueError:
                continue
        else:
            try:
                slide_numbers.append(int(part))
            except ValueError:
                continue

    # ソートして重複を除去、max_slides制限を適用
    result = sorted(set(slide_numbers))
    return result[:max_slides]


# =============================================================================
# Core Functions
# =============================================================================

def get_presentation_info(content: bytes, filename: str) -> PresentationInfo:
    """
    PowerPointプレゼンテーションの構造情報を取得する。

    Args:
        content: PowerPointファイルのバイトデータ
        filename: ファイル名

    Returns:
        PresentationInfo: プレゼンテーション情報
    """
    prs = _load_presentation_from_bytes(content)

    slides: list[SlideInfo] = []
    total_characters = 0

    for i, slide in enumerate(prs.slides, 1):
        title = _get_slide_title(slide)

        text_count = 0
        image_count = 0
        table_count = 0
        chart_count = 0
        char_count = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_count += 1
                for paragraph in shape.text_frame.paragraphs:
                    char_count += len(normalize_text(paragraph.text))
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1
            if shape.has_table:
                table_count += 1
                for row in shape.table.rows:
                    for cell in row.cells:
                        char_count += len(normalize_text(cell.text))
            if shape.has_chart:
                chart_count += 1

        # ノート
        has_notes = False
        notes_length = 0
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = normalize_text(slide.notes_slide.notes_text_frame.text).strip()
            if notes_text:
                has_notes = True
                notes_length = len(notes_text)
                char_count += notes_length

        total_characters += char_count

        slides.append(SlideInfo(
            number=i,
            title=title,
            text_count=text_count,
            image_count=image_count,
            table_count=table_count,
            chart_count=chart_count,
            char_count=char_count,
            has_notes=has_notes,
            notes_length=notes_length,
        ))

    return PresentationInfo(
        filename=filename,
        total_slides=len(prs.slides),
        total_characters=total_characters,
        slides=slides,
    )


def get_slides_content(
    content: bytes,
    *,
    slides_spec: str = "1-10",
    max_slides: int = DEFAULT_MAX_SLIDES,
    include_notes: bool = True,
    include_tables: bool = True,
) -> SlidesResult:
    """
    指定スライドの内容を取得する。

    Args:
        content: PowerPointファイルのバイトデータ
        slides_spec: スライド指定（"1-5" または "1,3,5" 形式）
        max_slides: 最大取得スライド数（デフォルト: 10）
        include_notes: ノートを含めるか（デフォルト: True）
        include_tables: 表を含めるか（デフォルト: True）

    Returns:
        SlidesResult: 取得結果
    """
    prs = _load_presentation_from_bytes(content)
    total_slides = len(prs.slides)

    slide_numbers = _parse_slides(slides_spec, total_slides, max_slides)

    slides: list[SlideContent] = []

    for slide_num in slide_numbers:
        if slide_num < 1 or slide_num > total_slides:
            continue

        slide = prs.slides[slide_num - 1]
        title = _get_slide_title(slide)

        # テキストを抽出
        text_content: list[str] = []
        table_content: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = normalize_text(paragraph.text).strip()
                    if text:
                        text_content.append(text)

            if include_tables and shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        row_texts.append(normalize_text(cell.text).strip())
                    if any(row_texts):
                        table_content.append("| " + " | ".join(row_texts) + " |")

        # ノート
        notes = None
        if include_notes and slide.has_notes_slide:
            notes_text = normalize_text(slide.notes_slide.notes_text_frame.text).strip()
            if notes_text:
                notes = notes_text

        slides.append(SlideContent(
            number=slide_num,
            title=title,
            text_content=text_content,
            table_content=table_content,
            notes=notes,
        ))

    # 範囲情報を計算
    if slide_numbers:
        start_slide = min(slide_numbers)
        end_slide = max(slide_numbers)
    else:
        start_slide = 0
        end_slide = 0

    has_more = end_slide < total_slides

    return SlidesResult(
        filename="",  # ハンドラーで設定
        requested_slides=slides_spec,
        total_slides=total_slides,
        returned_slides=len(slides),
        start_slide=start_slide,
        end_slide=end_slide,
        has_more=has_more,
        slides=slides,
    )


def search_presentation(
    content: bytes,
    query: str,
    *,
    case_sensitive: bool = False,
    max_hits: int = 50,
    include_notes: bool = True,
) -> SearchResult:
    """
    プレゼンテーション全体からキーワード検索を行う。

    Args:
        content: PowerPointファイルのバイトデータ
        query: 検索キーワード
        case_sensitive: 大文字小文字を区別するか（デフォルト: False）
        max_hits: 最大ヒット数（デフォルト: 50）
        include_notes: ノートも検索対象に含めるか（デフォルト: True）

    Returns:
        SearchResult: 検索結果
    """
    prs = _load_presentation_from_bytes(content)

    hits: list[SearchHit] = []
    pattern = build_search_pattern(query, case_sensitive)

    for slide_num, slide in enumerate(prs.slides, 1):
        if len(hits) >= max_hits:
            break

        title = _get_slide_title(slide)

        # タイトルを検索
        match = pattern.search(title)
        if match and len(hits) < max_hits:
            hits.append(SearchHit(
                slide_number=slide_num,
                slide_title=title,
                location_type="title",
                text=title,
                context=create_context_snippet(title, match.start(), match.end()),
            ))

        # テキストを検索
        for shape in slide.shapes:
            if len(hits) >= max_hits:
                break

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if len(hits) >= max_hits:
                        break

                    text = normalize_text(paragraph.text)
                    if not text:
                        continue

                    match = pattern.search(text)
                    if match:
                        hits.append(SearchHit(
                            slide_number=slide_num,
                            slide_title=title,
                            location_type="text",
                            text=text[:200] if len(text) > 200 else text,
                            context=create_context_snippet(text, match.start(), match.end()),
                        ))

            # 表を検索
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    if len(hits) >= max_hits:
                        break

                    row_texts = []
                    for cell in row.cells:
                        cell_text = normalize_text(cell.text).strip()
                        row_texts.append(cell_text)

                        match = pattern.search(cell_text)
                        if match and len(hits) < max_hits:
                            row_context = " | ".join(row_texts)
                            hits.append(SearchHit(
                                slide_number=slide_num,
                                slide_title=title,
                                location_type="table",
                                text=cell_text[:200] if len(cell_text) > 200 else cell_text,
                                context=f"行{row_idx + 1}: {row_context[:100]}",
                            ))
                            break  # 1行につき1ヒットまで

        # ノートを検索
        if include_notes and slide.has_notes_slide and len(hits) < max_hits:
            notes_text = normalize_text(slide.notes_slide.notes_text_frame.text).strip()
            if notes_text:
                match = pattern.search(notes_text)
                if match:
                    hits.append(SearchHit(
                        slide_number=slide_num,
                        slide_title=title,
                        location_type="notes",
                        text=notes_text[:200] if len(notes_text) > 200 else notes_text,
                        context=create_context_snippet(notes_text, match.start(), match.end()),
                    ))

    return SearchResult(
        query=query,
        total_hits=len(hits),
        hits=hits,
    )


# =============================================================================
# Tool Handlers
# =============================================================================

@file_tool_handler(
    old_format=(".ppt", "PowerPoint", ".pptx", "python-pptx", "Microsoft PowerPoint"),
    required_library=("pptx", "python-pptx"),
    log_prefix="PowerPoint情報取得",
)
async def get_presentation_info_handler(*, content, filename, args, **_):
    """
    PowerPointプレゼンテーションの構造情報を取得するハンドラー

    Args:
        args:
            file_path: ファイルパス
    """
    info = get_presentation_info(content, filename)

    result_lines = [
        f"# PowerPoint情報: {info['filename']}",
        f"スライド数: {info['total_slides']}",
        f"総文字数: {info['total_characters']:,}",
        "",
        "## スライド一覧",
    ]

    for slide in info['slides']:
        elements = []
        if slide['text_count'] > 0:
            elements.append(f"テキスト{slide['text_count']}")
        if slide['image_count'] > 0:
            elements.append(f"画像{slide['image_count']}")
        if slide['table_count'] > 0:
            elements.append(f"表{slide['table_count']}")
        if slide['chart_count'] > 0:
            elements.append(f"グラフ{slide['chart_count']}")

        element_str = ", ".join(elements) if elements else "空"

        result_lines.append(f"")
        result_lines.append(f"### スライド {slide['number']} - \"{slide['title']}\"")
        result_lines.append(f"- 要素: {element_str}")
        result_lines.append(f"- 文字数: 約{slide['char_count']}文字")
        if slide['has_notes']:
            result_lines.append(f"- ノート: あり ({slide['notes_length']}文字)")

    result_lines.append("")
    result_lines.append("---")
    result_lines.append("スライド取得: `get_slides_content` を使用")
    result_lines.append("検索: `search_presentation` を使用")

    return format_tool_success("\n".join(result_lines))


@file_tool_handler(
    old_format=(".ppt", "PowerPoint", ".pptx", "python-pptx", "Microsoft PowerPoint"),
    required_library=("pptx", "python-pptx"),
    log_prefix="PowerPointスライド取得",
)
async def get_slides_content_handler(*, content, filename, args, **_):
    """
    PowerPointスライドの内容を取得するハンドラー

    Args:
        args:
            file_path: ファイルパス
            slides: スライド指定（例: "1-5" または "1,3,5"）
            max_slides: 最大取得スライド数（デフォルト: 10）
            include_notes: ノートを含めるか（デフォルト: true）
            include_tables: 表を含めるか（デフォルト: true）
    """
    slides_spec = args.get("slides", "1-10")
    max_slides = args.get("max_slides", DEFAULT_MAX_SLIDES)
    include_notes = args.get("include_notes", True)
    include_tables = args.get("include_tables", True)

    result = get_slides_content(
        content,
        slides_spec=slides_spec,
        max_slides=max_slides,
        include_notes=include_notes,
        include_tables=include_tables,
    )

    result_lines = [
        f"# {filename}",
        f"取得スライド: {result['requested_slides']} (全{result['total_slides']}スライド)",
        f"返却: {result['returned_slides']}スライド",
        "",
    ]

    for slide in result['slides']:
        result_lines.append(f"## スライド {slide['number']} - \"{slide['title']}\"")
        result_lines.append("")

        if slide['text_content']:
            for text in slide['text_content']:
                result_lines.append(text)
            result_lines.append("")

        if slide['table_content']:
            result_lines.append("### 表")
            for row in slide['table_content']:
                result_lines.append(row)
            result_lines.append("")

        if slide['notes']:
            result_lines.append("### ノート")
            result_lines.append(slide['notes'])
            result_lines.append("")

        if not slide['text_content'] and not slide['table_content']:
            result_lines.append("[このスライドにテキストは含まれていません]")
            result_lines.append("")

        result_lines.append("---")
        result_lines.append("")

    if result['has_more']:
        result_lines.append("まだ続きがあります。次を取得するには:")
        result_lines.append(f"`slides=\"{result['end_slide'] + 1}-{result['end_slide'] + max_slides}\"` を指定してください。")

    return format_tool_success("\n".join(result_lines))


@file_tool_handler(
    old_format=(".ppt", "PowerPoint", ".pptx", "python-pptx", "Microsoft PowerPoint"),
    required_library=("pptx", "python-pptx"),
    log_prefix="PowerPoint検索",
)
async def search_presentation_handler(*, content, args, **_):
    """
    PowerPointプレゼンテーション全体からキーワード検索を行うハンドラー

    Args:
        args:
            file_path: ファイルパス
            query: 検索キーワード
            case_sensitive: 大文字小文字を区別するか（デフォルト: false）
            max_hits: 最大ヒット数（デフォルト: 50）
            include_notes: ノートも検索対象に含めるか（デフォルト: true）
    """
    query = args.get("query", "")
    case_sensitive = args.get("case_sensitive", False)
    max_hits = args.get("max_hits", 50)
    include_notes = args.get("include_notes", True)

    if not query:
        return format_tool_error("エラー: query（検索キーワード）を指定してください。")

    result = search_presentation(
        content,
        query,
        case_sensitive=case_sensitive,
        max_hits=max_hits,
        include_notes=include_notes,
    )

    result_lines = [
        f"# 検索結果: \"{result['query']}\"",
        f"ヒット数: {result['total_hits']}",
        "",
    ]

    if result['hits']:
        for hit in result['hits']:
            location_label = {
                "title": "タイトル",
                "text": "テキスト",
                "table": "表",
                "notes": "ノート",
            }.get(hit['location_type'], hit['location_type'])

            result_lines.append(f"## スライド {hit['slide_number']} - \"{hit['slide_title']}\"")
            result_lines.append(f"場所: {location_label}")
            result_lines.append(f"コンテキスト: {hit['context']}")
            result_lines.append("")
    else:
        result_lines.append("検索結果はありませんでした。")

    return format_tool_success("\n".join(result_lines))
