"""
ファイル読み込みMCPツール

Lazy Loading: AIエージェントがオンデマンドでファイルを読み込む

ツール一覧:
- read_image_file: 画像ファイルを読み込み（Base64 content block）
- read_pdf_file: PDFファイルを読み込み（document block）
- read_office_file: Office系ファイル（Excel/Word/PowerPoint）を読み込み
- list_workspace_files: ワークスペース内のファイル一覧を取得
"""

import base64
import io
from pathlib import Path
from typing import TYPE_CHECKING, Any, Optional

import structlog

from app.services.workspace.file_processors import FileCategory, FileTypeClassifier

if TYPE_CHECKING:
    from app.services.workspace_service import WorkspaceService

logger = structlog.get_logger(__name__)


async def _extract_excel_content(
    content: bytes, sheet_name: Optional[str] = None, max_rows: int = 1000
) -> str:
    """
    Excelファイルからテキストを抽出

    Args:
        content: ファイル内容
        sheet_name: 読み込むシート名（省略時は全シート）
        max_rows: 最大行数

    Returns:
        抽出されたテキスト
    """
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "Error: openpyxlライブラリがインストールされていません。"

    try:
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

        result = []
        sheets = [sheet_name] if sheet_name else wb.sheetnames

        for sheet in sheets:
            if sheet not in wb.sheetnames:
                continue
            ws = wb[sheet]
            result.append(f"## シート: {sheet}\n")

            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows:
                    result.append(f"\n(以降 {ws.max_row - max_rows} 行省略)\n")
                    break
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                rows.append(row_text)

            result.append("\n".join(rows))
            result.append("\n")

        wb.close()
        return "\n".join(result)
    except Exception as e:
        logger.error("Excel読み込みエラー", error=str(e))
        return f"Error: Excelファイルの読み込みに失敗しました: {str(e)}"


async def _extract_docx_content(content: bytes) -> str:
    """
    Wordファイルからテキストを抽出

    Args:
        content: ファイル内容

    Returns:
        抽出されたテキスト
    """
    try:
        from docx import Document
    except ImportError:
        return "Error: python-docxライブラリがインストールされていません。"

    try:
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.error("Word読み込みエラー", error=str(e))
        return f"Error: Wordファイルの読み込みに失敗しました: {str(e)}"


async def _extract_pptx_content(content: bytes) -> str:
    """
    PowerPointファイルからテキストを抽出

    Args:
        content: ファイル内容

    Returns:
        抽出されたテキスト
    """
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: python-pptxライブラリがインストールされていません。"

    try:
        prs = Presentation(io.BytesIO(content))

        result = []
        for i, slide in enumerate(prs.slides, 1):
            result.append(f"## スライド {i}\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    result.append(shape.text)
            result.append("")

        return "\n".join(result)
    except Exception as e:
        logger.error("PowerPoint読み込みエラー", error=str(e))
        return f"Error: PowerPointファイルの読み込みに失敗しました: {str(e)}"


def create_file_reader_handlers(
    workspace_service: "WorkspaceService",
    tenant_id: str,
    conversation_id: str,
):
    """
    ファイル読み込みハンドラーを作成

    Args:
        workspace_service: WorkspaceServiceインスタンス
        tenant_id: テナントID
        conversation_id: 会話ID

    Returns:
        ハンドラー辞書
    """

    async def read_image_file_handler(args: dict[str, Any]) -> dict[str, Any]:
        """画像ファイル読み込み"""
        file_path = args.get("file_path", "")

        try:
            content, filename, content_type = await workspace_service.download_file(
                tenant_id, conversation_id, file_path
            )

            # サポートされている画像形式かチェック
            category = FileTypeClassifier.get_category(filename, content_type)
            if category != FileCategory.IMAGE:
                return {
                    "content": [
                        {
                            "type": "text",
                            "text": f"このファイルは画像ではありません: {filename} ({content_type})",
                        }
                    ],
                    "is_error": True,
                }

            # Base64エンコード
            base64_data = base64.b64encode(content).decode("utf-8")

            # Messages API用のcontent block形式で返す
            return {
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": content_type,
                            "data": base64_data,
                        },
                    }
                ],
                "_metadata": {
                    "file_path": file_path,
                    "file_size": len(content),
                    "content_type": content_type,
                },
            }
        except FileNotFoundError:
            return {
                "content": [{"type": "text", "text": f"ファイルが見つかりません: {file_path}"}],
                "is_error": True,
            }
        except Exception as e:
            logger.error("画像ファイル読み込みエラー", error=str(e), file_path=file_path)
            return {
                "content": [{"type": "text", "text": f"読み込みエラー: {str(e)}"}],
                "is_error": True,
            }

    async def read_pdf_file_handler(args: dict[str, Any]) -> dict[str, Any]:
        """PDFファイル読み込み"""
        file_path = args.get("file_path", "")

        try:
            content, filename, content_type = await workspace_service.download_file(
                tenant_id, conversation_id, file_path
            )

            # PDFファイルかチェック
            category = FileTypeClassifier.get_category(filename, content_type)
            if category != FileCategory.PDF:
                return {
                    "content": [
                        {
                            "type": "text",
                            "text": f"このファイルはPDFではありません: {filename} ({content_type})",
                        }
                    ],
                    "is_error": True,
                }

            # Base64エンコード
            base64_data = base64.b64encode(content).decode("utf-8")

            # Messages API用のdocument block形式で返す
            return {
                "content": [
                    {
                        "type": "document",
                        "source": {
                            "type": "base64",
                            "media_type": "application/pdf",
                            "data": base64_data,
                        },
                    }
                ],
                "_metadata": {
                    "file_path": file_path,
                    "file_size": len(content),
                },
            }
        except FileNotFoundError:
            return {
                "content": [{"type": "text", "text": f"ファイルが見つかりません: {file_path}"}],
                "is_error": True,
            }
        except Exception as e:
            logger.error("PDFファイル読み込みエラー", error=str(e), file_path=file_path)
            return {
                "content": [{"type": "text", "text": f"読み込みエラー: {str(e)}"}],
                "is_error": True,
            }

    async def read_office_file_handler(args: dict[str, Any]) -> dict[str, Any]:
        """Officeファイル読み込み（テキスト抽出）"""
        file_path = args.get("file_path", "")
        sheet_name = args.get("sheet_name")
        max_rows = args.get("max_rows", 1000)

        # ファイル拡張子で処理を分岐
        ext = Path(file_path).suffix.lower()

        try:
            content, filename, content_type = await workspace_service.download_file(
                tenant_id, conversation_id, file_path
            )

            # Officeファイルかチェック
            category = FileTypeClassifier.get_category(filename, content_type)
            if category != FileCategory.OFFICE:
                return {
                    "content": [
                        {
                            "type": "text",
                            "text": f"このファイルはOffice形式ではありません: {filename} ({content_type})",
                        }
                    ],
                    "is_error": True,
                }

            # 古いOffice形式（.doc/.xls/.ppt）はサポートされていない
            # python-docx / python-pptx は Open XML 形式（.docx/.pptx）のみ対応
            # openpyxl は .xlsx のみ対応（.xls は xlrd が必要だが現在は未対応）
            if ext == ".doc":
                return {
                    "content": [{
                        "type": "text",
                        "text": (
                            f"古いWord形式（.doc）はサポートされていません: {filename}\n"
                            ".docx形式に変換してから再度アップロードしてください。"
                        ),
                    }],
                    "is_error": True,
                }
            elif ext == ".ppt":
                return {
                    "content": [{
                        "type": "text",
                        "text": (
                            f"古いPowerPoint形式（.ppt）はサポートされていません: {filename}\n"
                            ".pptx形式に変換してから再度アップロードしてください。"
                        ),
                    }],
                    "is_error": True,
                }
            elif ext == ".xls":
                return {
                    "content": [{
                        "type": "text",
                        "text": (
                            f"古いExcel形式（.xls）はサポートされていません: {filename}\n"
                            ".xlsx形式に変換してから再度アップロードしてください。"
                        ),
                    }],
                    "is_error": True,
                }
            elif ext == ".xlsx":
                # Excelファイルの処理
                text_content = await _extract_excel_content(content, sheet_name, max_rows)
            elif ext == ".docx":
                # Wordファイルの処理
                text_content = await _extract_docx_content(content)
            elif ext == ".pptx":
                # PowerPointファイルの処理
                text_content = await _extract_pptx_content(content)
            else:
                return {
                    "content": [{"type": "text", "text": f"未対応のOffice形式: {ext}"}],
                    "is_error": True,
                }

            return {
                "content": [{"type": "text", "text": text_content}],
                "_metadata": {
                    "file_path": file_path,
                    "file_size": len(content),
                    "extracted_length": len(text_content),
                },
            }
        except FileNotFoundError:
            return {
                "content": [{"type": "text", "text": f"ファイルが見つかりません: {file_path}"}],
                "is_error": True,
            }
        except Exception as e:
            logger.error("Officeファイル読み込みエラー", error=str(e), file_path=file_path)
            return {
                "content": [{"type": "text", "text": f"読み込みエラー: {str(e)}"}],
                "is_error": True,
            }

    async def list_workspace_files_handler(args: dict[str, Any]) -> dict[str, Any]:
        """ワークスペースファイル一覧取得"""
        filter_type = args.get("filter_type", "all")

        try:
            file_list = await workspace_service.list_files(tenant_id, conversation_id)

            files_info = []
            for f in file_list.files:
                category = FileTypeClassifier.get_category(f.file_path, f.mime_type)
                category_str = category.value

                if filter_type != "all" and category_str != filter_type:
                    continue

                files_info.append(
                    {
                        "path": f.file_path,
                        "name": f.original_name,
                        "size": f.file_size,
                        "type": category_str,
                        "mime_type": f.mime_type,
                    }
                )

            result_text = f"ワークスペース内のファイル一覧（{len(files_info)}件）:\n\n"
            for f in files_info:
                size_kb = f["size"] / 1024
                result_text += f"- {f['path']} ({size_kb:.1f}KB, {f['type']})\n"

            if not files_info:
                result_text = "ワークスペースにファイルがありません。"

            return {
                "content": [{"type": "text", "text": result_text}],
                "_metadata": {"file_count": len(files_info), "files": files_info},
            }
        except Exception as e:
            logger.error("ファイル一覧取得エラー", error=str(e))
            return {
                "content": [{"type": "text", "text": f"エラー: {str(e)}"}],
                "is_error": True,
            }

    return {
        "read_image_file": read_image_file_handler,
        "read_pdf_file": read_pdf_file_handler,
        "read_office_file": read_office_file_handler,
        "list_workspace_files": list_workspace_files_handler,
    }


# ファイル読み込みツールに関するシステムプロンプト
FILE_READER_PROMPT = """
## ファイル読み込みツール

ワークスペースにアップロードされたファイルを読み込むには、以下のツールを使用してください：

### ファイル一覧の確認
- `mcp__file-reader__list_workspace_files`: ワークスペース内のファイル一覧を取得
  - filter_type: "image" / "pdf" / "office" / "text" / "all"

### 画像ファイルの読み込み
- `mcp__file-reader__read_image_file`: 画像ファイル（JPEG/PNG/GIF/WebP）を読み込み
  - file_path: ファイルパス（例: "uploads/photo.jpg"）
  - 画像の内容を視覚的に分析できます

### PDFファイルの読み込み
- `mcp__file-reader__read_pdf_file`: PDFファイルを読み込み
  - file_path: ファイルパス（例: "uploads/document.pdf"）
  - テキストと画像の両方を含むPDFを処理できます

### Officeファイルの読み込み
- `mcp__file-reader__read_office_file`: Excel/Word/PowerPointファイルを読み込み
  - file_path: ファイルパス（例: "uploads/data.xlsx"）
  - sheet_name: Excelの場合、読み込むシート名（省略時は全シート）
  - max_rows: Excelの場合、読み込む最大行数（デフォルト: 1000）

### 重要な注意
- ファイルを分析する前に、必ず `list_workspace_files` で存在を確認してください
- 大きなファイルは読み込みに時間がかかる場合があります
- テキスト/CSV/JSONファイルは従来の `Read` ツールも使用可能です
"""
